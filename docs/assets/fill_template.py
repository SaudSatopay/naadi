"""Fill the official IDBI Innovate 2026 template with NAADI content.

Reads  : docs/assets/idbi-template.pptx  + console screenshots in shots/
Writes : docs/NAADI-IDBI-submission.pptx
Body areas are empty in the template; we add textboxes/images/shapes and
leave every branded header, footer and cover element untouched.
"""

from __future__ import annotations

import copy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
SHOTS = HERE / "shots"
OUT = HERE.parent / "NAADI-IDBI-submission.pptx"

INK = RGBColor(0x20, 0x2B, 0x27)
MUT = RGBColor(0x5B, 0x6B, 0x64)
GRN = RGBColor(0x0E, 0x9F, 0x6E)
ORG = RGBColor(0xEA, 0x58, 0x0C)
RED = RGBColor(0xC2, 0x41, 0x0C)
CARD = RGBColor(0xF2, 0xF7, 0xF4)
LINE = RGBColor(0xD5, 0xE2, 0xDB)

BODY_TOP = Inches(1.45)
BODY_LEFT = Inches(0.35)
BODY_W = Inches(9.3)
BODY_BOT = Inches(5.35)


# --------------------------------------------------------------- helpers ----
def crop(src: str, box: tuple[int, int, int, int], name: str) -> Path:
    out = SHOTS / f"crop-{name}.png"
    Image.open(SHOTS / src).crop(box).save(out)
    return out


def tb(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size=11, bold=False, color=INK, first=False, align=None,
         space_after=6, bullet=False, font="Arial"):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    if align:
        p.alignment = align
    p.space_after = Pt(space_after)
    chunks = text.split("**")
    for i, chunk in enumerate(chunks):
        if not chunk:
            continue
        r = p.add_run()
        r.text = ("• " if bullet and i == 0 else "") + chunk
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold or (i % 2 == 1)
        r.font.color.rgb = color
    return p


def card(slide, x, y, w, h, fill=CARD):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = LINE
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def flow_box(slide, x, y, w, h, title, sub, accent=GRN):
    sh = card(slide, x, y, w, h)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.name = "Arial"; r.font.size = Pt(10.5); r.font.bold = True
    r.font.color.rgb = accent
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sub
    r2.font.name = "Arial"; r2.font.size = Pt(8); r2.font.color.rgb = INK
    return sh


def arrow(slide, x, y, w=Inches(0.22)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.16))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x9C, 0xB5, 0xAA)
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def pic(slide, path: Path, x, y, w=None, h=None):
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def stat(slide, x, y, w, big, label, color=GRN):
    tf = tb(slide, x, y, w, Inches(1.1))
    para(tf, big, size=30, bold=True, color=color, first=True, align=PP_ALIGN.CENTER, space_after=2)
    para(tf, label, size=9, color=MUT, align=PP_ALIGN.CENTER)


# ----------------------------------------------------------------- crops ----
c_hero = crop("book-tall.png", (200, 70, 1400, 560), "hero")
c_map = crop("book-tall.png", (200, 560, 1400, 1010), "map")
c_cards = crop("book-tall.png", (200, 1040, 1400, 1500), "cards")
c_dims = crop("chandra-tall.png", (700, 680, 1370, 1020), "dims")
c_traj = crop("chandra-tall.png", (230, 1030, 1370, 1300), "traj")
c_water = crop("chandra-tall.png", (230, 1470, 1370, 1745), "water")
c_reason = crop("chandra-tall.png", (230, 1750, 1370, 2140), "reason")
c_sens = crop("chandra-tall.png", (230, 2305, 1370, 2600), "sens")
c_memo = crop("chandra-tall.png", (230, 2600, 1370, 3300), "memo")

prs = Presentation(HERE / "idbi-template.pptx")
S = list(prs.slides)


# ---------------------------------------------------------- slide 1: team ----
for p in S[0].shapes[1].text_frame.paragraphs:
    joined = "".join(r.text for r in p.runs)
    if "Team name" in joined and p.runs:
        p.runs[0].text = "Team name:  NAADI"
    elif "Team leader" in joined and p.runs:
        p.runs[0].text = "Team leader name:  Saud Satopay"
    elif "Problem Statement" in joined and p.runs:
        p.runs[0].text = ("Problem Statement:  Track 03 — Financial Health Score "
                          "(AI/ML MSME Financial Health Card on GST · UPI · AA · EPFO)")

# ---------------------------------------------------------- slide 2: brief ----
s = S[1]
tf = tb(s, BODY_LEFT, BODY_TOP, Inches(4.6), Inches(3.8))
para(tf, "NAADI (नाड़ी — “the pulse”)", size=17, bold=True, color=GRN, first=True, space_after=8)
para(tf, "An AI-native **MSME Financial Health Card** plus **Munshi**, an agentic "
         "underwriting copilot. NAADI fuses four consented data rails — **GST, UPI, "
         "Account Aggregator, EPFO** — into a six-dimension, explainable 300–900 score, "
         "a credit decision with a limit, and a committee-ready sanction memo.", size=11.5, space_after=10)
para(tf, "score → reasons → decision → memo → monitoring", size=12, bold=True, color=ORG, space_after=10)
para(tf, "Built for **New-to-Credit / New-to-Bank** enterprises that have no audited "
         "financials and no bureau file — the exact gap the Bank's problem statement names. "
         "Working prototype: trained model (held-out **AUROC 0.871**), live console, "
         "decision API — all public.", size=11.5)
pic(s, c_hero, Inches(5.15), BODY_TOP, w=Inches(4.5))

# --------------------------------------------------- slide 3: opportunities ----
s = S[2]
# blank the template's prompt bullets (keep the "Opportunities" heading)
for sh in s.shapes:
    if sh.has_text_frame and "Opportunities" in sh.text_frame.text:
        for p in sh.text_frame.paragraphs[1:]:
            for r in p.runs:
                r.text = ""
cols = [
    ("How is it different?", GRN, [
        "**Six explainable dimensions** + honest confidence band — not one opaque number",
        "**Monotone scorecards + TreeSHAP** reason codes a credit committee can audit",
        "**Thin-file tiers** widen uncertainty instead of rejecting data-poor borrowers",
        "**LLM writes prose, never numbers** — deterministic engine owns every figure",
    ]),
    ("How does it solve the problem?", ORG, [
        "Scores **credit-invisible MSMEs in seconds** from consented alternate data",
        "Cuts application-to-decision **from days to minutes** (Munshi drafts the memo)",
        "**Rolling 12-month rescore** flags stress quarters early — origination + monitoring",
        "Every decline ships a **what-if comeback path** (e.g. +24 pts for on-time GST ×3 mo)",
    ]),
    ("USP", GRN, [
        "**An underwriting brain, not a score** — decision, limit, covenants, triggers included",
        "**Interrogable AI**: sensitivity sliders answered by the trained model itself",
        "**Human gate built in**: officer concurs or overrides on the record",
        "**ULI / OCEN / AA-native** design with a working POST /score API today",
    ]),
]
cw = Inches(3.0)
for i, (title, accent, items) in enumerate(cols):
    x = BODY_LEFT + Emu(int((cw + Inches(0.15)) * i))
    card(s, x, BODY_TOP, cw, Inches(3.75))
    tf = tb(s, x + Inches(0.15), BODY_TOP + Inches(0.12), cw - Inches(0.3), Inches(3.5))
    para(tf, title, size=13, bold=True, color=accent, first=True, space_after=8)
    for it in items:
        para(tf, it, size=10, bullet=True, space_after=7)

# -------------------------------------------------------- slide 4: features ----
s = S[3]
feats = [
    "**Six-dimension Health Card** — Liquidity · Stability · Growth · Repayment · Compliance · Concentration",
    "**Calibrated PD(12m) → 300–900 score** with tier-based confidence band",
    "**TreeSHAP reason codes** in plain language (+ score waterfall)",
    "**What-if engine** — actionable counterfactual paths for the borrower",
    "**Sensitivity lab** — drag a lever, the trained model projects the score",
    "**Stress lab** — adverse scenarios rescored, not guessed",
    "**Munshi memo** — decision, limit, tenor, covenants, EWS triggers (Claude-drafted, engine-numbered)",
    "**Officer review gate** — concur / override with recorded rationale",
    "**Portfolio radar** — rolling rescore + rule-based EWS alert feed",
    "**Risk map** — whole book by score × PD × exposure",
    "**⇄ Compare** — two files on one radar with dimension deltas",
    "**Decision API** — POST /score for LOS / ULI / OCEN integration",
]
half = 6
for col in range(2):
    x = BODY_LEFT + Emu(int(Inches(4.75) * col))
    tf = tb(s, x, BODY_TOP, Inches(4.55), Inches(3.9))
    for i, f in enumerate(feats[col * half:(col + 1) * half]):
        para(tf, f, size=10.5, bullet=True, first=(i == 0), space_after=9)

# ------------------------------------------------------- slide 5: process ----
# Serpentine flow: row 1 left->right, wrap down on the right, row 2 right->left.
# Each step carries three detail lines; outcomes + feedback loop below.
s = S[4]


def flow_detail(slide, x, y, w, h, title, lines, accent=GRN):
    sh = card(slide, x, y, w, h)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(3)
    r = p.add_run(); r.text = title
    r.font.name = "Arial"; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = accent
    for ln in lines:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_after = Pt(1)
        r2 = p2.add_run(); r2.text = ln
        r2.font.name = "Arial"; r2.font.size = Pt(7.3); r2.font.color.rgb = INK
    return sh


def arrow_left(slide, x, y, w=Inches(0.2)):
    a = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, x, y, w, Inches(0.16))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x9C, 0xB5, 0xAA)
    a.line.fill.background(); a.shadow.inherit = False
    return a


def arrow_down(slide, x, y, h=Inches(0.28)):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, Inches(0.16), h)
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x9C, 0xB5, 0xAA)
    a.line.fill.background(); a.shadow.inherit = False
    return a


row1 = [
    ("1 · CONSENT", ["AA consent artefact", "GSTN OTP · DPDP purpose-bound", "immutable audit ledger"]),
    ("2 · PULL RAILS", ["GSTR-1/3B returns", "UPI txns · AA bank statements", "EPFO payroll (ECR)"]),
    ("3 · FEATURES", ["27 point-in-time signals", "leak-proof as-of joins", "tier T1/T2/T3 coverage"]),
    ("4 · SIX SCORECARDS", ["monotone LightGBM", "L·S·G·R·C·K subscores 0–100", "sector-peer percentile scaled"]),
    ("5 · SCORE + PD", ["composite → isotonic PD(12m)", "300–900 with ± tier band", "benchmarked vs 5,000 peers"]),
]
row2_right_to_left = [
    ("6 · EXPLAIN", ["TreeSHAP reason codes", "what-if comeback paths", "sensitivity + stress labs"]),
    ("7 · POLICY + LIMIT", ["cash-flow affordability (EMI room)", "guards: L-floor · K-cap · C-gate", "limit · tenor · covenants"]),
    ("8 · MUNSHI MEMO", ["Claude drafts the prose only", "every figure injected by engine", "evidence-linked, committee-ready"]),
    ("9 · OFFICER GATE", ["a human always decides", "concur / override + rationale", "recorded to audit trail"]),
    ("10 · MONITOR", ["rolling monthly rescore", "EWS alerts below 580", "portfolio radar feed"]),
]

bw, bh, gap = Inches(1.68), Inches(1.06), Inches(0.20)
y1 = BODY_TOP + Inches(0.06)
y2 = y1 + bh + Inches(0.52)
xs = [BODY_LEFT + Emu(int((bw + gap) * i)) for i in range(5)]

for i, (t, lines) in enumerate(row1):
    flow_box_accent = GRN
    flow_detail(s, xs[i], y1, bw, bh, t, lines, flow_box_accent)
    if i < 4:
        arrow(s, xs[i] + bw + Inches(0.0), y1 + Inches(0.45))

# wrap-down connector on the right edge: step 5 -> step 6
arrow_down(s, xs[4] + Emu(int(bw / 2)) - Inches(0.08), y1 + bh + Inches(0.10), h=Inches(0.32))

# row 2 runs right -> left (serpentine), so step 6 sits under step 5
for j, (t, lines) in enumerate(row2_right_to_left):
    i = 4 - j  # column index from the right
    accent = ORG if t.split(" · ")[1] in ("MUNSHI MEMO", "OFFICER GATE") else GRN
    flow_detail(s, xs[i], y2, bw, bh, t, lines, accent)
    if j < 4:
        arrow_left(s, xs[i] - Inches(0.20), y2 + Inches(0.45))

# outcomes strip under the officer gate + monitor
oy = y2 + bh + Inches(0.16)
outcomes = [
    ("✓  APPROVE — sanction with limit, tenor, covenants", GRN),
    ("◐  REFER — committee with mitigants", ORG),
    ("✕  DECLINE — with a what-if comeback plan", RED),
]
ow = Inches(2.95)
for i, (t, c) in enumerate(outcomes):
    x = BODY_LEFT + Emu(int((ow + Inches(0.22)) * i))
    ch = card(s, x, oy, ow, Inches(0.34))
    tf = ch.text_frame
    tf.margin_left = Inches(0.08); tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t
    r.font.name = "Arial"; r.font.size = Pt(8.5); r.font.bold = True; r.font.color.rgb = c

# feedback loop: MONITOR (bottom-left) feeds fresh data back into FEATURES
from pptx.enum.shapes import MSO_CONNECTOR  # noqa: E402
conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                              xs[0] + Emu(int(bw / 2)), y2 - Inches(0.02),
                              xs[2] + Emu(int(bw / 2)), y1 + bh + Inches(0.02))
conn.line.color.rgb = RGBColor(0x9C, 0xB5, 0xAA)
conn.line.width = Pt(1.1)
conn.line.dash_style = 2  # dashed
tf = tb(s, xs[0] + Inches(0.5), y2 - Inches(0.34), Inches(2.6), Inches(0.28))
para(tf, "feedback: monthly rescore", size=8, color=MUT, first=True)

tf = tb(s, BODY_LEFT, oy + Inches(0.44), BODY_W, Inches(0.32))
para(tf, "Every run is consent-first and fully audited; the officer — not the model — owns the final decision.",
     size=9.5, color=MUT, first=True, align=PP_ALIGN.CENTER)

# ----------------------------------------------------- slide 6: wireframes ----
s = S[5]
pic(s, c_map, BODY_LEFT, BODY_TOP, w=Inches(5.55))
pic(s, c_dims, Inches(6.1), BODY_TOP, w=Inches(3.35))
tf = tb(s, Inches(6.1), BODY_TOP + Inches(1.95), Inches(3.35), Inches(1.8))
para(tf, "Live console (not a mock):", size=11, bold=True, color=GRN, first=True, space_after=6)
para(tf, "risk map of the book + EWS alert feed (left) · the six-dimension "
         "card grid (above) — all rendered from engine output.", size=10, color=INK)
para(tf, "naadi-kappa.vercel.app", size=11, bold=True, color=ORG)

# --------------------------------------------------- slide 7: architecture ----
s = S[6]
lanes = [
    ("CONSENTED RAILS", GRN, [
        "GSTN — GSTR-1/3B, filing discipline",
        "UPI — daily pulse, payer graph",
        "Account Aggregator — balances, bounces",
        "EPFO — headcount, wage trend",
        "Consent ledger — DPDP, immutable audit",
    ]),
    ("THE TWO-LAYER BRAIN", GRN, [
        "27 point-in-time features (Polars, leak-proof)",
        "6 monotone LightGBM scorecards (L·S·G·R·C·K)",
        "Composite → isotonic PD(12m) → 300–900 ± band",
        "TreeSHAP reasons · what-if · sensitivity · stress",
        "Peer benchmarks vs 5,000-MSME population",
    ]),
    ("MUNSHI + SURFACES", ORG, [
        "Deterministic policy & limit engine",
        "Claude-drafted memo (prose only, fact-sheet injected)",
        "Officer gate: concur / override, recorded",
        "RM console · print-ready sanction memo",
        "POST /score API → LOS · ULI · OCEN rails",
    ]),
]
lw = Inches(2.92)
for i, (title, accent, items) in enumerate(lanes):
    x = BODY_LEFT + Emu(int((lw + Inches(0.42)) * i))
    card(s, x, BODY_TOP, lw, Inches(3.45))
    tf = tb(s, x + Inches(0.15), BODY_TOP + Inches(0.12), lw - Inches(0.3), Inches(3.2))
    para(tf, title, size=12, bold=True, color=accent, first=True, space_after=7)
    for it in items:
        para(tf, it, size=9.5, bullet=True, space_after=6)
    if i < 2:
        arrow(s, x + lw + Inches(0.07), BODY_TOP + Inches(1.55), w=Inches(0.28))
tf = tb(s, BODY_LEFT, BODY_TOP + Inches(3.6), BODY_W, Inches(0.35))
para(tf, "Demo: Polars + DuckDB on laptop  →  Scale path: Kafka → Iceberg · Feast · MLflow · EKS · OpenTelemetry — "
         "connectors swap to IDBI sandbox with zero downstream changes.",
     size=9.5, color=MUT, first=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------- slide 8: tech ----
s = S[7]
tech = [
    ("Console", "Next.js 16 (App Router, Turbopack) · React 19 · Tailwind CSS v4 · Motion 12 · Recharts 3 — fully static"),
    ("Engine", "Python 3.13 · uv · Polars · DuckDB · scikit-learn 1.9 · FastAPI (POST /score, p95 < 300 ms)"),
    ("ML", "LightGBM 4.6 monotone scorecards · isotonic calibration · native TreeSHAP · counterfactual & sensitivity engines"),
    ("Copilot", "Claude Opus 4.8 (adaptive thinking) for memo prose · deterministic offline fallback — demo can never break"),
    ("Data & consent", "Synthetic latent-risk generator today → IDBI sandbox APIs (Jul 22) · AA/GSTN/EPFO connectors · DPDP-aligned consent ledger"),
    ("Scale path", "AWS (EKS) · Kafka → Iceberg lakehouse · Feast feature store · MLflow registry · OpenTelemetry + Langfuse"),
]
for i, (k, v) in enumerate(tech):
    y = BODY_TOP + Emu(int(Inches(0.63) * i))
    tf = tb(s, BODY_LEFT, y, Inches(1.7), Inches(0.55))
    para(tf, k, size=12, bold=True, color=GRN if i % 2 == 0 else ORG, first=True)
    tf = tb(s, Inches(2.15), y, Inches(7.5), Inches(0.55))
    para(tf, v, size=10.5, first=True)

# ---------------------------------------------------------- slide 9: cost ----
s = S[8]
rows = [
    ("Concept build (done)", "₹0 infrastructure — static console (Vercel), model trains on a laptop in ~40 s", GRN),
    ("Sandbox PoC (Jul 22–31)", "Covered by hackathon-provided AWS credits + sandbox APIs; no bank spend", GRN),
    ("Pilot (indicative)", "Serving: containerised FastAPI on 2 small nodes ≈ ₹25–40k/mo · memo LLM ≈ ₹1–3 per file (offline fallback → ₹0)", ORG),
    ("At scale (indicative)", "< ₹2 marginal compute per decision at 10k decisions/day; data-rail costs per GSP/AA pulls as contracted", ORG),
]
for i, (k, v, c) in enumerate(rows):
    y = BODY_TOP + Emu(int(Inches(0.78) * i))
    card(s, BODY_LEFT, y, BODY_W, Inches(0.68))
    tf = tb(s, BODY_LEFT + Inches(0.2), y + Inches(0.09), Inches(2.4), Inches(0.5))
    para(tf, k, size=11, bold=True, color=c, first=True)
    tf = tb(s, Inches(3.1), y + Inches(0.09), Inches(6.35), Inches(0.5))
    para(tf, v, size=10, first=True)
tf = tb(s, BODY_LEFT, BODY_TOP + Inches(3.3), BODY_W, Inches(0.4))
para(tf, "Cost discipline is architectural: the LLM only writes prose (small, cacheable calls); scoring is pure CPU.",
     size=10, color=MUT, first=True, align=PP_ALIGN.CENTER)

# ------------------------------------------------------ slide 10: snapshots ----
s = S[9]
pic(s, c_water, BODY_LEFT, BODY_TOP, w=Inches(4.85))
pic(s, c_reason, BODY_LEFT, BODY_TOP + Inches(1.35), w=Inches(4.85))
pic(s, c_memo, Inches(5.45), BODY_TOP, h=Inches(3.75))
tf = tb(s, BODY_LEFT, BODY_TOP + Inches(3.15), Inches(4.85), Inches(0.6))
para(tf, "Score waterfall · TreeSHAP reason ledger · stress lab (left) — Munshi's "
         "evidence-linked sanction memo (right). All live at naadi-kappa.vercel.app.",
     size=9.5, color=MUT, first=True)

# ---------------------------------------------------- slide 11: benchmarks ----
s = S[10]
stat(s, BODY_LEFT, BODY_TOP, Inches(2.3), "0.871", "held-out AUROC — composite PD model")
stat(s, Inches(2.75), BODY_TOP, Inches(2.3), "0.576", "KS statistic (test fold)")
stat(s, Inches(5.15), BODY_TOP, Inches(2.3), "5,000", "synthetic MSME validation population", ORG)
stat(s, Inches(7.4), BODY_TOP, Inches(2.2), "< 300 ms", "p95 · POST /score decision API", ORG)
pic(s, c_traj, BODY_LEFT, BODY_TOP + Inches(1.35), w=Inches(5.9))
tf = tb(s, Inches(6.5), BODY_TOP + Inches(1.4), Inches(3.1), Inches(2.4))
para(tf, "Early-warning proof:", size=11, bold=True, color=RED, first=True, space_after=5)
para(tf, "rolling rescore kept the failing supplier **below the EWS line for all 12 months** "
         "before its PD hit 63%.", size=10, space_after=8)
para(tf, "Why not “90% accuracy”? Default is a rare event — accuracy is a broken yardstick. "
         "We report **AUROC / KS + calibration** and will defend that choice.", size=9.5, color=MUT)
tf = tb(s, BODY_LEFT, BODY_TOP + Inches(3.55), BODY_W, Inches(0.35))
para(tf, "Synthetic-population validation of the architecture — not a market claim. Re-estimated on IDBI sandbox data (Jul 22–31) with a published report.",
     size=9.5, color=MUT, first=True, align=PP_ALIGN.CENTER)

# -------------------------------------------------------- slide 12: future ----
s = S[11]
future = [
    ("Jul 22–31 · Sandbox", "Retrain + isotonic recalibration on IDBI data; publish validation report (AUROC/KS/PSI, fairness slices); live-data demo"),
    ("PoC · LOS integration", "POST /score into loan origination; swap-set study — how many good NTC borrowers were being declined?"),
    ("Rails", "ULI / OCEN 4.0 embedded journeys — score at the point of commerce; AA-consent UX hardening"),
    ("Borrower side", "Vernacular MSME app: my Health Card, my what-if plan, my 90-day “build your NAADI” track"),
    ("Monitoring", "Productise the radar into a Track-04-grade early-warning system for the whole MSME book"),
]
for i, (k, v) in enumerate(future):
    y = BODY_TOP + Emu(int(Inches(0.72) * i))
    tf = tb(s, BODY_LEFT, y, Inches(2.5), Inches(0.6))
    para(tf, k, size=11.5, bold=True, color=GRN if i < 2 else ORG, first=True)
    tf = tb(s, Inches(3.0), y, Inches(6.6), Inches(0.6))
    para(tf, v, size=10.5, first=True)

# --------------------------------------------------------- slide 13: links ----
s = S[12]
# the template textbox lists prompt bullets where our content goes — keep the
# heading line, blank the bullets
for sh in s.shapes:
    if sh.has_text_frame and "Provide links" in sh.text_frame.text:
        for p in sh.text_frame.paragraphs[1:]:
            for r in p.runs:
                r.text = ""
rows = [
    ("GitHub (public)", "github.com/SaudSatopay/naadi", "qr-repo.png"),
    ("Final product (live)", "naadi-kappa.vercel.app", "qr-live.png"),
]
for i, (k, v, qr) in enumerate(rows):
    y = BODY_TOP + Inches(0.25) + Emu(int(Inches(1.55) * i))
    pic(s, HERE / qr, BODY_LEFT, y, w=Inches(1.15))
    tf = tb(s, Inches(1.75), y + Inches(0.25), Inches(6.6), Inches(0.9))
    para(tf, k, size=12, bold=True, color=MUT, first=True, space_after=3)
    para(tf, v, size=18, bold=True, color=GRN)
tf = tb(s, BODY_LEFT, BODY_TOP + Inches(3.45), Inches(9.0), Inches(0.4))
para(tf, "In-repo documentation: ARCHITECTURE.md · SCORING.md · PITCH.md · 9-page concept deck (PDF)",
     size=10, color=MUT, first=True)

# ------------------------------------------------------- slide 14: closer ----
s = S[13]
tf = tb(s, Inches(0.6), Inches(2.1), Inches(8.8), Inches(1.6))
para(tf, "score → reasons → decision → memo → monitoring", size=16, bold=True, color=GRN,
     first=True, align=PP_ALIGN.CENTER, space_after=10)
para(tf, "NAADI — because the businesses Bharat runs on deserve to be seen.",
     size=14, color=INK, align=PP_ALIGN.CENTER)

prs.save(OUT)
print("saved", OUT)
